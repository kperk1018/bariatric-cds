"""Build the deck (PowerPoint) — slides + clinician-facing speaker notes.

Audience: a surgeon / clinical team with limited machine-learning background.
Every ML term is explained in plain language, and every slide's notes end on the
CLINICAL SIGNIFICANCE — why it matters for a patient. The speaker notes are the
script: read them and you can present without knowing ML.

Run:
    PYTHONPATH=. python scripts/build_presentation.py
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from src.config import BASELINE_FEATURES

FIG = Path("presentation/figures")
OUT = Path("presentation/Bariatric_CDS_Raftopoulos_2026-07-13.pptx")

NAVY = RGBColor(0x1F, 0x35, 0x4E)
GREY = RGBColor(0x54, 0x6E, 0x7A)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1E, 0x8A, 0x4C)
W, H = Inches(13.333), Inches(7.5)


def _title(s, text, sub=None):
    tb = s.shapes.add_textbox(Inches(0.45), Inches(0.24), Inches(12.5), Inches(0.72))
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY
    if sub:
        tb2 = s.shapes.add_textbox(Inches(0.45), Inches(0.94), Inches(12.5), Inches(0.4))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = sub; p2.font.size = Pt(13); p2.font.color.rgb = GREY; p2.font.italic = True


def _bul(s, items, x=0.45, y=1.42, w=12.5, h=4.0, size=14):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        text, bold, color = (it if isinstance(it, tuple) else (it, False, None))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.bold = bold
        p.font.color.rgb = color or RGBColor(0x21, 0x21, 0x21)
        p.space_after = Pt(6)


def _pic(s, name, x, y, w):
    f = FIG / name
    if f.exists():
        s.shapes.add_picture(str(f), Inches(x), Inches(y), width=Inches(w))


def _notes(s, t):
    s.notes_slide.notes_text_frame.text = t


def build():
    prs = Presentation(); prs.slide_width, prs.slide_height = W, H
    B = prs.slide_layouts[6]

    # 1 ── Title
    s = prs.slides.add_slide(B)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "Predicting Weight-Loss Trajectories After Bariatric Surgery"
    p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = NAVY
    tb2 = s.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.7))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Flagging the patients likely to regain — before they do"
    p2.font.size = Pt(18); p2.font.color.rgb = GREY
    _bul(s, [("786 patients  •  a personalized forecast  •  5 trajectory types  •  a working clinical tool",
              True, NAVY)], y=4.25, size=15)
    _notes(s, """WHAT TO SAY:
"Every patient asks two questions before surgery: how much weight will I lose, and will I keep it
off? Today, the honest answer is a population average — 'most people lose about 30%.' But patients
vary enormously, and some regain. And usually, by the time we SEE the regain in clinic, the window
to intervene has already closed."

"This tool does three things: it gives a forecast for THIS patient rather than the average patient;
it tells you how much to trust that forecast, and refuses to answer when it can't; and it sorts
patients into recognizable trajectory types so we can spot the likely regainers early."

CLINICAL SIGNIFICANCE: the whole point is earlier, more targeted follow-up — spending our limited
clinic time on the patients who actually need it, before they've regained.""")

    # 2 ── Attrition Sankey
    s = prs.slides.add_slide(B)
    _title(s, "First — who do we actually still have?",
           "Real-world follow-up drops off fast; this sets what the tool can and can't tell you")
    _pic(s, "fig0b_sankey.png", 0.35, 1.55, 12.6)
    _notes(s, """WHAT TO SAY:
"Before any accuracy number, look at who we still have. We start with 786 patients. By year 1, 593
have follow-up data. By year 2, 386. By year 5, only 85 — about 10% of the cohort. People stop
coming to clinic. This is the reality of long-term bariatric follow-up."

"This one picture explains everything that follows: why our year 2-4 predictions are strong, and
why years 5 and 6 are unreliable. It's not that the tool is bad late — it's that there's almost
nothing left to learn from."

A REASSURING POINT TO MAKE: the patients who DO stay to years 5-6 tend to be the ones who started
with LESS weight loss before surgery — the harder cases. So if anything, the late numbers understate
how well people do, rather than exaggerate it.

CLINICAL SIGNIFICANCE: this is about setting honest expectations. The tool is genuinely useful
through year 4 — the window where you can still act — and it's upfront about the limits beyond that,
rather than giving you a confident number you shouldn't trust.""")

    # 3 ── Pipeline + the 15 fields
    s = prs.slides.add_slide(B)
    _title(s, "How it works — start to finish")
    _pic(s, "fig1_pipeline.png", 0.3, 1.35, 12.7)
    feats = ", ".join(BASELINE_FEATURES)
    _bul(s, [
        ("THE 15 THINGS THE TOOL LOOKS AT — all collected routinely before surgery:", True, NAVY),
        (feats, False, None),
        ("It cleans the records, removes duplicate entries, learns the patterns ONCE, and then simply "
         "looks up an answer for each new patient. Nothing is re-computed per visit — that's why it's "
         "instant.", False, None),
    ], y=4.35, size=12)
    _notes(s, """WHAT TO SAY — walk the boxes left to right:
"Raw records come in. We clean them and remove duplicate rows — there were 16 duplicated patient
records, which we caught and removed, leaving 786 real patients. Then each patient is reduced to 15
preoperative numbers."

NAME THE 15 FIELDS out loud (they're on the slide): age, sex, race, height, initial BMI, initial
weight, initial metabolic rate, visceral fat, body-fat %, fat mass, fat-free mass, time to surgery,
surgery type, preop BMI, and preop weight loss. Emphasize: "every one of these is something you
already collect — the tool needs no new tests."

"Then the prediction engine runs, produces the forecast plus a trust rating, groups the patient into
a trajectory type, and shows it all in the app."

IF ASKED 'does it re-learn each time?': No — it's trained once and saved, like writing down a recipe.
Each patient just gets cooked from the recipe. That's why it responds instantly.

CLINICAL SIGNIFICANCE: it runs entirely on information already in the preoperative workup — no extra
burden on the patient or the clinic.""")

    # 4 ── What a Random Forest is
    s = prs.slides.add_slide(B)
    _title(s, "The prediction engine, in plain terms",
           "Hundreds of simple rules, voting")
    _pic(s, "fig2_model_comparison.png", 0.3, 1.5, 8.6)
    _bul(s, [
        ("THE ANALOGY", True, NAVY),
        "Imagine asking 300 junior doctors to each guess a patient's weight loss. Each is only allowed "
        "a few yes/no questions — Is BMI over 45? Did they lose 10% before surgery? Each guess is "
        "mediocre. But AVERAGE all 300, and the errors cancel out.",
        "That's the model. Each 'junior doctor' is a simple decision tree; hundreds of them vote.",
        ("TWO SAFEGUARDS", True, NAVY),
        "It never peeks at the future: the year-2 forecast only uses information available by year 1.",
        "It's reproducible: run it twice, get identical numbers.",
    ], x=9.0, y=1.55, w=4.1, h=5.4, size=11)
    _notes(s, """WHAT TO SAY — define it with the analogy, don't assume the audience knows it:
"The engine is called a Random Forest, which is an unhelpful name. Here's what it actually is.
Imagine 300 junior doctors, each allowed only a handful of yes/no questions about the patient. Each
one alone gives a so-so guess. But average all 300 guesses together and the random errors cancel out,
leaving something surprisingly accurate. Each 'junior doctor' is what's called a decision tree, and
hundreds of them vote. No black box — just many simple rules, averaged."

THE TWO SAFEGUARDS (a skeptical clinician will wonder about these):
1. No cheating with future information. If the year-2 forecast were allowed to peek at the year-2
   answer, it would look brilliant and be useless in real life. We checked every model — the year-2
   forecast only ever sees information available by year 1.
2. Reproducible. There's randomness inside, so we fix it in place; anyone re-running gets the exact
   same numbers.

CLINICAL SIGNIFICANCE: this is an interpretable, checkable method — not a mysterious AI. It behaves
consistently, and it can't secretly rely on information you wouldn't have at the bedside.""")

    # 5 ── RF vs GB
    s = prs.slides.add_slide(B)
    _title(s, "We tested two versions and kept the steadier one",
           "Peak accuracy in one year matters less than reliability across all years")
    _pic(s, "fig2b_rf_vs_gb.png", 0.5, 1.4, 10.4)
    _bul(s, [
        ("Random Forest (used) — R²:  yr1 0.25 · yr2 0.51 · yr3 0.61 · yr4 0.73", True, NAVY),
        ("Random Forest (used) — AUC: yr1 0.74 · yr2 0.90 · yr3 0.91 · yr4 0.95", True, NAVY),
        ("Gradient Boosting edges ahead at year 3, but fades in the sparse late years; "
         "Random Forest stays steadier across all years, so the app uses it and keeps GB as a "
         "cross-check.", False, None),
    ], x=0.5, y=5.55, w=12.3, size=12)
    _notes(s, """WHAT TO SAY:
"There are two closely related versions of this kind of model. We ran both, head to head, on our own
patients. The second one — Gradient Boosting — is a bit sharper in the middle years; it actually
beats the first at year 3. But look at the late years: it collapses and becomes wildly unreliable,
while the first version degrades gently and predictably."

"So which do you want driving a clinical tool? Not the one that's occasionally sharper and
occasionally unhinged — you want the one that's steady and honest across every year. That's the one
the app uses. We keep the sharper version as a cross-check in the background."

IF ASKED 'why is the second one better at year 3-4?': it's a more aggressive learner, so when there's
enough data (the middle years) it squeezes out a bit more signal. But that same aggressiveness makes
it fall apart when data is thin (the late years) — it starts chasing noise.

CLINICAL SIGNIFICANCE: reliability is a safety property. A tool that's brilliant sometimes and badly
wrong other times is worse than one that's consistently good and knows its limits.""")

    # 6 ── Reliability + AUC
    s = prs.slides.add_slide(B)
    _title(s, "The tool refuses to guess when it can't predict",
           "Arguably the most important safety feature")
    _pic(s, "fig3_reliability.png", 0.6, 1.55, 12.1)
    _bul(s, [
        ("Two report cards per year. R² = can it predict the actual weight-loss NUMBER? AUC = can it "
         "at least tell a high responder from a low one? Green = trustworthy, amber = rough guide, "
         "red = the tool shows nothing at all.", True, RED),
        ("Sometimes it can't nail the exact % but can still reliably RANK who's at risk — that's still "
         "clinically useful, for triage rather than a precise number.", False, None),
    ], y=5.35, size=12.5)
    _notes(s, """WHAT TO SAY — define both report cards plainly:
"To decide whether to trust a prediction, we use two measures. R-squared asks: of all the variation
between patients, how much can the model actually explain? Zero means it knows nothing; higher is
better; above about 0.4 is genuinely useful in medicine. AUC asks a simpler question: can it at least
correctly rank patients — put the big losers above the small losers? Ours are 0.90, 0.91, 0.95 at
years 2, 3, 4, which is strong."

"We set the trust thresholds in advance. Green means trustworthy — years 2 to 4. Amber means treat as
a rough guide. Red means the model genuinely can't predict, and there the tool shows you NOTHING."

"That refusal is the single most important design choice. Most tools always give an answer — which is
dangerous, because a confident wrong number is worse than no number."

THE COMBINED VIEW (point to the blue box): "Sometimes the model can't predict the exact percentage
but can still reliably separate high from low responders — for example fat-loss at year 3: it can't
give you the number, but it can tell you which side of the line a patient is on. That's still useful
— for triage, not for a precise figure."

WHY THE MISSING HALF: even our best explains about half the variation. The rest is things we never
measured — behavior, mental health, life events, the support a patient has at home. Much of it likely
exists in scanned clinic notes that were never entered as data. So the gap isn't mysterious; it's
unmeasured, and it's fixable over time.

CLINICAL SIGNIFICANCE: the tool tells you when to trust it and when to fall back on your own judgment.
That honesty is what makes it safe to put in front of a patient.""")

    # 7 ── k / clustering
    s = prs.slides.add_slide(B)
    _title(s, "Finding the patient 'types' — grouped by their weight-loss journey",
           "The computer groups patients purely by the shape of their predicted curve")
    _pic(s, "fig4_silhouette.png", 0.35, 1.55, 6.4)
    _pic(s, "fig5_umap.png", 7.05, 1.55, 5.1)
    _bul(s, [
        ("Grouping means letting the computer find patients with similar weight-loss journeys — using "
         "ONLY the predicted trajectory, not their demographics or surgery type. We swept every group "
         "count from 2 to 10 and use five trajectory types, matching the established analysis.", False, None),
        ("The map on the right just squashes each patient down to a dot so you can see the groups "
         "separate — it's a picture, not the analysis.", False, None),
    ], y=6.0, size=12),
    _notes(s, """WHAT TO SAY — define the terms as you go:
"These patient types come purely from the SHAPE of the weight-loss journey — we do NOT feed in age,
sex, race, or surgery type. The computer groups patients only by their predicted curves. Anything we
later notice about who's in each group is a discovery, not something we told it to find."

"How many groups? We swept from 2 up to 10 and looked at a separation score. Mathematically the
cleanest split is just 2 big groups, but that's too coarse to be clinically useful — so we use five
trajectory types, matching the established analysis and giving the recognizable ladder from worst to
best responders. We're transparent that five is a clinically-informed choice, shown on the curve."

"The picture on the right just compresses each patient to a dot on a 2-D map so you can see the groups
pull apart. It's a visualization aid, not the analysis itself."

CLINICAL SIGNIFICANCE: because the groups come from the trajectories alone, they capture how a patient
actually does over time — not a stereotype based on their chart.""")

    # 8 ── The five trajectory types (descriptive composition)
    s = prs.slides.add_slide(B)
    _title(s, "The five trajectory types",
           "Groups come from the curves; who's in each is described afterward")
    _pic(s, "fig6_trajectories.png", 0.25, 1.45, 6.7)
    _pic(s, "fig7_demographics.png", 7.1, 1.45, 5.9)
    _notes(s, """WHAT TO SAY:
"Here are the five types. On the left, their weight-loss journeys, ordered from the worst responders
to the best. On the right, a description of who ended up in each group — remember, none of this was an
input; it's what we found after grouping by trajectory alone."

"Reading the ladder: the lowest group starts below the 10.5% preop weight-loss line — the level we
already know predicts worse long-term outcomes — and stays lowest. The best group are the gastric
bypass patients, who lose the most and hold it."

WHAT THE COMPOSITION SHOWS (describe, don't over-claim): the groups line up mainly with preoperative
weight loss and procedure. Race is mixed across the groups — it is NOT what separates them. Keep this
descriptive; do not frame it as 'the clusters are demographics.'

BE READY FOR ODD LATE NUMBERS: a couple of year-6 points rest on as few as 3 patients (labeled on the
figure) — quirks of tiny late samples, which is why years 5-6 are marked unreliable.

CLINICAL SIGNIFICANCE: recognizable patterns you can act on — the group below the 10.5% line is the
clearest early target for extra preoperative support.

(Then pivot: "One of these journeys hides a patient worth flagging — next slide.")""")

    # 9 ── The Quiet Regainer (the clinical payoff)
    s = prs.slides.add_slide(B)
    _title(s, "The 'Quiet Regainer' — the patient you'd otherwise miss",
           "A trajectory type that looks like a success early, then quietly regains")
    _pic(s, "fig_quiet_regainer.png", 0.9, 1.5, 11.5)
    _bul(s, [
        ("One trajectory type loses strongly early — looks like your best patients at year 1 — then "
         "quietly regains by year 4. You can't see it from the chart or the one-year visit; by the time "
         "it shows, the regain has happened.", True, RED),
        ("The tool flags this pattern UP FRONT from the predicted curve — buying two to three years of "
         "lead time to intervene. It's a built-in high-alert flag in the app.", True, GREEN),
    ], y=5.55, size=12),
    _notes(s, """*** THE CLINICAL PAYOFF OF THE WHOLE TALK ***

WHAT TO SAY:
"Among the trajectory types, one is clinically special. This group loses strongly early — at year 1
they look like your best patients, right up there with the top responders. But then they quietly
regain, sliding back down by year 4."

"This is the patient you would never flag from their chart, and you can't catch it at the one-year
visit because they look like a star. By the time the regain is visible, it's already happened."

"Because the tool works from the SHAPE of the predicted curve, it can flag these patients up front —
at the preop visit — giving you two or three years of lead time to intervene. It's a built-in
high-alert flag in the app."

HOW IT'S DEFINED (if asked): a simple, transparent rule on the predicted curve — a strong early peak
(≥34% by year 1–2) followed by a meaningful regain (≥12 percentage points down by year 4). No black
box; you can see exactly why a patient is flagged.

WHY ONLY THE TRAJECTORY VIEW FINDS THEM: two such patients can look identical before surgery and at
year 1 — they only diverge later. You can only separate them by looking at the whole journey.

CLINICAL SIGNIFICANCE: this is the single most actionable output of the tool — an early warning for
the exact patient who otherwise slips through until it's too late.""")

    # 10 ── Mechanism + next
    s = prs.slides.add_slide(B)
    _title(s, "Where a prediction comes from — and how you'd use it")
    _bul(s, [
        ("IT IS NOT A STRAIGHT LINE DRAWN FORWARD. IT IS NOT THE GROUP AVERAGE.", True, RED),
        "We take this individual patient's 15 preoperative values, run them through the model for that "
        "specific year, and hundreds of simple rules vote. The result is that patient's own predicted "
        "weight loss — each year computed separately.",
        "The order matters: we predict the trajectory FIRST, then read off which type it matches. The "
        "group doesn't create the number; the number places the patient in a group.",
        ("HOW YOU'D USE IT IN CLINIC", True, NAVY),
        "At the preoperative visit: a realistic, personalized expectation, plus a risk flag (below the "
        "10.5% threshold, or a likely Quiet Regainer).",
        "It sharpens further once you enter the patient's real year-1 result — so it's a follow-up tool "
        "as much as a counseling tool.",
        ("The working app is the deliverable — enter a patient, see their trajectory and their flags.",
         True, GREEN),
    ], y=1.45, size=13)
    _notes(s, """WHAT TO SAY — kill the two natural misreadings, because either one makes the tool sound
worthless:
"It is NOT extrapolation — we're not drawing a line through the early points and extending it. And it
is NOT the group average — we're not just telling you what that type of patient usually does."

"What actually happens: we take THIS patient's 15 preoperative values, run them through the model for
that specific year, hundreds of simple rules vote, and out comes that patient's own predicted weight
loss. Each year is computed independently."

"And the order matters: we predict the patient's trajectory first, and only then read off which of the
five types it matches. The group doesn't create the number — the number is what sorts the patient into
a group."

HOW IT'S USED: "At the preop visit, you get a realistic personalized expectation for counseling, plus
a risk flag. And it gets sharper the moment you enter a real follow-up measurement — so it's a
follow-up tool as much as a counseling tool."

CLINICAL SIGNIFICANCE / CLOSE: "The deliverable isn't a paper figure — it's a working tool. You enter
a patient, and you immediately see their likely trajectory and whether they're a high-alert Quiet
Regainer. That's the point: turning data we already collect into an early warning we can act on." """)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Deck written: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
